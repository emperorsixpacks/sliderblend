import os
import json
from pptx import Presentation

# Directory containing your PowerPoint templates
template_dir = "templates"  # Change this to your directory (use "." for current directory)

all_templates = {}

# Iterate over all .pptx files in the directory
for filename in os.listdir(template_dir):
    if filename.lower().endswith(".pptx"):
        template_path = os.path.join(template_dir, filename)
        prs = Presentation(template_path)
        layout_data = []

        # Extract layouts and their placeholders
        for i, layout in enumerate(prs.slide_layouts):
            placeholders = []
            for placeholder in layout.placeholders:
                placeholders.append({
                    "idx": placeholder.placeholder_format.idx,
                    "name": placeholder.name
                })
            layout_data.append({
                "index": i,
                "name": layout.name,
                "placeholders": placeholders
            })

        # Remove extension from filename for key
        template_key = os.path.splitext(filename)[0]
        
        # Store layout info along with the location of the template
        all_templates[template_key] = {
            "location": os.path.abspath(template_path),
            "layouts": layout_data
        }

# Write the aggregated data to a single JSON file named "templates.json" in the same directory
json_path = os.path.join(template_dir, "templates.json")
with open(json_path, "w") as f:
    json.dump(all_templates, f, indent=4)

print(f"Generated '{json_path}' with layouts and placeholders for all templates!")

